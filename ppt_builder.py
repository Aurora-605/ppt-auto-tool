from pptx import Presentation

def build_ppt(data):
    prs = Presentation()

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data["title"]
    slide.placeholders[1].text = "AI自动生成"

    # 内容页
    for slide_data in data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"]

        content = slide.placeholders[1].text_frame
        content.clear()

        for i, point in enumerate(slide_data["points"]):
            if i == 0:
                content.text = point
            else:
                p = content.add_paragraph()
                p.text = point

    output_file = "output.pptx"
    prs.save(output_file)

    return output_file
